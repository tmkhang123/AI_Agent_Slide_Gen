"""
slide_generator.py — "Full Bleed" layout
=========================================
Mỗi slide dùng ảnh làm nền toàn màn hình với hệ thống overlay thông minh:

  dark_bg   — ảnh tối  : overlay trắng rất mờ, gradient navy đậm ở dưới
  light_bg  — ảnh sáng : overlay đen vừa, gradient navy đậm ở dưới
  mid_bg    — trung bình: chỉ gradient bar phía dưới

  Fallback (không có ảnh): nền navy đặc

Title slide  : ảnh toàn màn + gradient panel trái → phải
Content slide: ảnh toàn màn + overlay + gradient bar phía dưới bảo vệ text
"""

from __future__ import annotations
from io import BytesIO

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE


NAVY     = RGBColor(0x1E, 0x27, 0x61)
ICE_BLUE = RGBColor(0xCA, 0xDC, 0xFC)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
MID_GRAY = RGBColor(0xCC, 0xCC, 0xDD)


# ─────────────────────────────────────────────────────────────────────────────
# Tiện ích xử lý ảnh (Pillow)
# ─────────────────────────────────────────────────────────────────────────────

def _brightness(stream: BytesIO) -> float:
    from PIL import Image, ImageStat
    stream.seek(0)
    img = Image.open(stream).convert("L").resize((64, 64))
    return ImageStat.Stat(img).mean[0]


def _prepare_bg(stream: BytesIO, w_in: float, h_in: float,
                blur: float = 0) -> BytesIO:
    """Crop-fill + optional blur. Trả về JPEG BytesIO."""
    from PIL import Image, ImageFilter
    DPI = 150
    tw, th = int(w_in * DPI), int(h_in * DPI)
    stream.seek(0)
    img = Image.open(stream).convert("RGB")
    iw, ih = img.size
    s = max(tw / iw, th / ih)
    img = img.resize((int(iw * s), int(ih * s)), Image.LANCZOS)
    nw, nh = img.size
    img = img.crop([(nw - tw) // 2, (nh - th) // 2,
                    (nw - tw) // 2 + tw, (nh - th) // 2 + th])
    if blur > 0:
        img = img.filter(ImageFilter.GaussianBlur(radius=blur))
    buf = BytesIO()
    img.save(buf, "JPEG", quality=88, optimize=True)
    buf.seek(0)
    return buf


def _solid_png(w_px: int, h_px: int, rgb: tuple, alpha: int) -> BytesIO:
    from PIL import Image
    buf = BytesIO()
    Image.new("RGBA", (w_px, h_px), (*rgb, alpha)).save(buf, "PNG")
    buf.seek(0)
    return buf


def _gradient_h(w_px: int, h_px: int, rgb: tuple,
                a_left: int, a_right: int) -> BytesIO:
    """Gradient ngang: alpha trái → alpha phải."""
    import numpy as np
    from PIL import Image
    row = np.linspace(a_left, a_right, w_px, dtype=np.uint8)
    a   = np.tile(row, (h_px, 1))
    r   = np.full_like(a, rgb[0])
    g   = np.full_like(a, rgb[1])
    b   = np.full_like(a, rgb[2])
    buf = BytesIO()
    Image.fromarray(np.stack([r, g, b, a], 2), "RGBA").save(buf, "PNG")
    buf.seek(0)
    return buf


def _gradient_v_bottom(w_px: int, h_px: int, rgb: tuple,
                       frac: float, a_top: int, a_bot: int) -> BytesIO:
    """Gradient dọc chỉ ở phần dưới: trong suốt → đặc."""
    import numpy as np
    from PIL import Image
    arr  = np.zeros((h_px, w_px, 4), dtype=np.uint8)
    bar  = int(h_px * frac)
    rows = np.linspace(a_top, a_bot, bar, dtype=np.uint8)
    for i, a in enumerate(rows):
        y = h_px - bar + i
        arr[y, :] = [rgb[0], rgb[1], rgb[2], a]
    buf = BytesIO()
    Image.fromarray(arr, "RGBA").save(buf, "PNG")
    buf.seek(0)
    return buf


# ─────────────────────────────────────────────────────────────────────────────
# SlideGenerator
# ─────────────────────────────────────────────────────────────────────────────

class SlideGenerator:
    W   = 10.0
    H   = 7.5
    DPI = 150

    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width  = Inches(self.W)
        self.prs.slide_height = Inches(self.H)

    # ── helpers ──────────────────────────────────────────────────────────────

    def _blank(self):
        for lay in self.prs.slide_layouts:
            if lay.name == "Blank":
                return self.prs.slides.add_slide(lay)
        return self.prs.slides.add_slide(self.prs.slide_layouts[6])

    def _solid_bg(self, slide, color):
        f = slide.background.fill
        f.solid(); f.fore_color.rgb = color

    def _rect(self, slide, l, t, w, h, rgb, border=False):
        sh = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(l), Inches(t), Inches(w), Inches(h))
        sh.fill.solid(); sh.fill.fore_color.rgb = rgb
        if border:
            sh.line.color.rgb = ICE_BLUE; sh.line.width = Pt(0.5)
        else:
            sh.line.fill.background()
        return sh

    def _pic(self, slide, stream: BytesIO, l, t, w, h):
        stream.seek(0)
        slide.shapes.add_picture(
            stream, Inches(l), Inches(t), width=Inches(w), height=Inches(h))

    def _tb(self, slide, l, t, w, h):
        tb = slide.shapes.add_textbox(
            Inches(l), Inches(t), Inches(w), Inches(h))
        tb.text_frame.word_wrap = True
        return tb

    def _run(self, para, text, pt, bold=False, italic=False,
             color=WHITE, font="Calibri"):
        r = para.add_run()
        r.text = text
        r.font.size = Pt(pt); r.font.bold = bold; r.font.italic = italic
        r.font.color.rgb = color; r.font.name = font
        return r

    def _layout_type(self, stream: BytesIO) -> str:
        try:
            b = _brightness(stream)
            if b < 85:   return "dark"
            if b > 155:  return "light"
            return "mid"
        except Exception:
            return "mid"

    # ── px helpers ───────────────────────────────────────────────────────────

    @property
    def _wp(self): return int(self.W * self.DPI)
    @property
    def _hp(self): return int(self.H * self.DPI)

    # ─────────────────────────────────────────────────────────────────────────
    # Title slide
    # ─────────────────────────────────────────────────────────────────────────

    def add_title_slide(self, title: str, subtitle: str,
                        image_stream: BytesIO | None = None):
        slide = self._blank()
        W, H  = self.W, self.H

        if image_stream is not None:
            # 1. Nền ảnh full-bleed (blur rất nhẹ)
            self._pic(slide, _prepare_bg(image_stream, W, H, blur=1.5),
                      0, 0, W, H)

            # 2. Gradient trái-phải: navy đặc bên trái → trong suốt bên phải
            #    Che 70% chiều rộng để bảo vệ vùng chữ
            PANEL_W = 0.72  # tỉ lệ chiều rộng slide
            self._pic(slide,
                      _gradient_h(int(self._wp * PANEL_W), self._hp,
                                  (0x1E, 0x27, 0x61), 235, 0),
                      0, 0, W * PANEL_W, H)
        else:
            self._solid_bg(slide, NAVY)

        # 3. Accent bar trái
        self._rect(slide, 0, 0, 0.18, H, ICE_BLUE)
        # 4. Gạch ngang phân cách
        self._rect(slide, 0.38, 2.9, 5.5, 0.045, ICE_BLUE)

        # 5. Tiêu đề
        tb = self._tb(slide, 0.55, 1.15, 5.7, 1.95)
        p  = tb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        self._run(p, title, 34, bold=True)

        # 6. Phụ đề
        tb2 = self._tb(slide, 0.55, 3.2, 5.7, 1.5)
        p2  = tb2.text_frame.paragraphs[0]
        p2.alignment = PP_ALIGN.LEFT
        self._run(p2, subtitle, 15, italic=True, color=ICE_BLUE)

        # 7. Watermark
        tb3 = self._tb(slide, 0.55, H - 0.55, 5.7, 0.35)
        p3  = tb3.text_frame.paragraphs[0]
        self._run(p3, "AI Slides Maker · Powered by Ollama", 9, color=MID_GRAY)

    # ─────────────────────────────────────────────────────────────────────────
    # Content slide
    # ─────────────────────────────────────────────────────────────────────────

    def add_content_slide(self, title: str, bullet_points: list,
                          image_stream: BytesIO | None = None):
        slide = self._blank()
        W, H  = self.W, self.H
        wp, hp = self._wp, self._hp
        HEADER_H = 1.06

        if image_stream is not None:
            ltype = self._layout_type(image_stream)

            # 1. Nền ảnh full-bleed (không blur để giữ chi tiết)
            self._pic(slide, _prepare_bg(image_stream, W, H, blur=0),
                      0, 0, W, H)

            # 2. Overlay theo độ sáng
            if ltype == "light":
                # Ảnh quá sáng: overlay đen vừa để chữ đọc được
                self._pic(slide, _solid_png(wp, hp, (0, 0, 0), 100),
                          0, 0, W, H)
            elif ltype == "dark":
                # Ảnh tối: overlay trắng rất nhạt, ảnh vẫn nổi
                self._pic(slide, _solid_png(wp, hp, (255, 255, 255), 18),
                          0, 0, W, H)
            # mid: không overlay toàn màn

            # 3. Gradient bar dọc phía dưới (bảo vệ vùng text chính)
            #    Từ trong suốt → navy đậm, chiếm 62% chiều cao từ dưới lên
            self._pic(slide,
                      _gradient_v_bottom(wp, hp, (0x1E, 0x27, 0x61),
                                         frac=0.62, a_top=0, a_bot=215),
                      0, 0, W, H)

            # 4. Header overlay tối để tiêu đề đọc rõ
            self._pic(slide, _solid_png(wp, int(HEADER_H * self.DPI),
                                        (0, 0, 0), 145),
                      0, 0, W, HEADER_H)
        else:
            self._solid_bg(slide, NAVY)

        # ── Header bar ──
        self._rect(slide, 0, 0, 0.18, HEADER_H, ICE_BLUE)
        self._rect(slide, 0, HEADER_H, W, 0.04, ICE_BLUE)

        tb_t = self._tb(slide, 0.38, 0.1, 9.2, 0.86)
        p_t  = tb_t.text_frame.paragraphs[0]
        p_t.alignment = PP_ALIGN.LEFT
        self._run(p_t, title, 22, bold=True)

        # ── Bullet text ──
        BT = HEADER_H + 0.12
        BH = H - BT - 0.28
        tb_b = self._tb(slide, 0.45, BT, 9.1, BH)
        tf   = tb_b.text_frame
        tf.word_wrap = True
        first = True
        for pt in bullet_points:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.space_before = Pt(8)
            p.space_after  = Pt(3)
            self._run(p, f"▪  {pt}", 15)

        # ── Footer ──
        tb_f = self._tb(slide, 0.38, H - 0.22, 9.2, 0.18)
        p_f  = tb_f.text_frame.paragraphs[0]
        p_f.alignment = PP_ALIGN.RIGHT
        self._run(p_f, title, 8, color=MID_GRAY)

    # ── Save ─────────────────────────────────────────────────────────────────

    def save(self, filename: str = "presentation.pptx") -> str:
        self.prs.save(filename)
        return filename
