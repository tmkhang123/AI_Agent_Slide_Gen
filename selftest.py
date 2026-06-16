"""
selftest.py — KIỂM TRA NHANH (không cần Ollama)
===============================================
Chạy:  python selftest.py

Xác minh:
  1. Bộ phân tích văn bản (parse_titles / parse_points) hoạt động đúng.
  2. Toàn bộ pipeline 3 agent dựng được file PPTX ở CHẾ ĐỘ DỰ PHÒNG
     (giả lập Ollama tắt) — bảo đảm phần "khung" luôn chạy được.

Nếu selftest PASS mà chạy thật vẫn ra dự phòng => vấn đề nằm ở Ollama
(chưa chạy 'ollama serve' hoặc chưa 'ollama pull <model>'),
hãy bấm nút "Kiểm tra Ollama" trong GUI.
"""

import llm
from agent_planner import PlannerAgent
from agent_content import ContentAgent
from agent_reviewer import ReviewerAgent
from agent_designer import DesignerAgent
import os
from pptx import Presentation


def test_parsers() -> bool:
    ok = True
    t = llm.parse_titles("Đây là danh sách:\n1. Giới thiệu :: intro concept\n"
                         "- Phân tích | analysis data\nSlide 3: Kết luận")
    ok &= len(t) == 3
    p = llm.parse_points("Nội dung:\n- Một ý dài đủ để vượt qua kiểm tra số từ "
                         "tối thiểu của hệ thống parser.\n"
                         "- Ý thứ hai cũng đủ dài và mang thông tin cụ thể rõ ràng.")
    ok &= len(p) == 2
    print(f"[1] Parsers: {'PASS' if ok else 'FAIL'} "
          f"(titles={len(t)}, points={len(p)})")
    return ok


def test_pipeline_fallback() -> bool:
    # Giả lập Ollama tắt: mọi lời gọi generate đều ném lỗi
    def down(*a, **k):
        raise ConnectionError("giả lập Ollama tắt")
    llm.generate = down

    outline = PlannerAgent().plan("Chủ đề kiểm thử")
    deck = ContentAgent().write(outline)
    reviewed_deck = ReviewerAgent().review(deck)
    pptx = DesignerAgent(use_images=False).build(reviewed_deck, "Product", "selftest")

    n = len(Presentation(pptx).slides._sldIdLst)
    has_bolds = any("**" in pt for s in reviewed_deck.slides for pt in s.points)
    
    ok = (len(reviewed_deck.slides) >= 6
          and all(len(s.points) >= 3 for s in reviewed_deck.slides)
          and n == len(reviewed_deck.slides) + 1
          and has_bolds
          and os.path.exists(pptx))
          
    print(f"[2] Pipeline (fallback): {'PASS' if ok else 'FAIL'} "
          f"({len(reviewed_deck.slides)} slide nội dung, PPTX {n} slide, bôi đậm={has_bolds})")
    return ok


def test_markdown_stripper() -> bool:
    from slide_generator import _strip_markdown
    ok = True
    ok &= _strip_markdown("**tổng quan về thuật toán Viterbi**") == "tổng quan về thuật toán Viterbi"
    ok &= _strip_markdown("<b>cơ chế phân tích cú pháp</b>") == "cơ chế phân tích cú pháp"
    ok &= _strip_markdown("__test__ and *another* test") == "test and another test"
    print(f"[3] Title Markdown Stripper: {'PASS' if ok else 'FAIL'}")
    return ok


def test_instruction_leaks() -> bool:
    text = (
        "- **Kiểm tra tính thực chất:** Đã kiểm tra xong.\n"
        "- **Bôi đậm từ khóa:** Đã làm.\n"
        "- Tối ưu độ dài: Rút gọn.\n"
        "- **Ý thứ nhất thực sự hữu ích** cho bài thuyết trình."
    )
    points = llm.parse_points(text)
    # Nó phải loại bỏ hết các dòng chỉ dẫn và chỉ giữ lại dòng nội dung thực sự.
    ok = len(points) == 1 and points[0] == "**Ý thứ nhất thực sự hữu ích** cho bài thuyết trình"
    print(f"[4] Instruction Leaks Filter: {'PASS' if ok else 'FAIL'} (parsed={points})")
    return ok


def test_output_delimiter() -> bool:
    text = (
        "Báo cáo phân tích:\n"
        "Tôi đã hoàn tất việc tối ưu slide này.\n"
        "- Lộ ý rác 1\n"
        "---OUTPUT---\n"
        "- Ý thứ nhất thực sự hữu ích\n"
        "- Ý thứ hai chuyên sâu kỹ thuật\n"
    )
    points = llm.parse_points(text)
    ok = len(points) == 2 and points[0] == "Ý thứ nhất thực sự hữu ích" and points[1] == "Ý thứ hai chuyên sâu kỹ thuật"
    print(f"[5] Output Delimiter: {'PASS' if ok else 'FAIL'} (parsed={points})")
    return ok


def test_title_formatting() -> bool:
    from agent_planner import _format_presentation_title
    ok = True
    # Test markdown stripping in formatting
    ok &= _format_presentation_title("**tổng quan về bert**") == "Tổng quan về BERT"
    ok &= _format_presentation_title("__vận hành của lstm__") == "Vận hành của LSTM"
    ok &= _format_presentation_title("rnn và cnn trong xử lý ảnh và ngôn ngữ") == "RNN và CNN trong xử lý ảnh và ngôn ngữ"
    print(f"[6] Title Formatting Acronyms: {'PASS' if ok else 'FAIL'}")
    return ok


def test_deep_leakage_filter() -> bool:
    text = (
        "- Loại bỏ số liệu fabricated và thuật ngữ lạ/sai thực chất\n"
        "- CẤM BỊA ĐẶT SỐ LIỆU & THUẬT NGỮ GIẢ: Tôi đã kiểm tra kỹ lưỡng.\n"
        "- **Ý thứ nhất** của slide thực sự học thuật và chính xác."
    )
    points = llm.parse_points(text)
    ok = len(points) == 1 and points[0] == "**Ý thứ nhất** của slide thực sự học thuật và chính xác"
    print(f"[7] Deep Leakage Filter: {'PASS' if ok else 'FAIL'} (parsed={points})")
    return ok


def test_deduplication() -> bool:
    text = (
        "- Ý thứ nhất học thuật\n"
        "- Ý thứ hai chuyên sâu\n"
        "- **Ý thứ nhất học thuật**\n"
        "- Ý thứ hai chuyên sâu\n"
        "- **Ý thứ hai chuyên sâu**"
    )
    points = llm.parse_points(text)
    # Nó phải khử trùng lặp và giữ lại phiên bản bôi đậm ở sau cùng
    ok = (len(points) == 2 
          and points[0] == "**Ý thứ nhất học thuật**" 
          and points[1] == "**Ý thứ hai chuyên sâu**")
    print(f"[8] Core Deduplication: {'PASS' if ok else 'FAIL'} (parsed={points})")
    return ok


def test_agenda_review_bypass() -> bool:
    from schema import Deck, Slide
    deck = Deck(
        title="Test Deck",
        subtitle="Subtitle",
        slides=[
            Slide(title="Chương trình nghị sự", points=["Phần 1: Giới thiệu", "Phần 2: Phân tích"]),
            Slide(title="Nội dung chính", points=["Ý thứ nhất cần tối ưu"])
        ]
    )
    orig_gen = llm.generate
    def down(*a, **k):
        raise ConnectionError("giả lập Ollama tắt")
    llm.generate = down
    try:
        reviewed = ReviewerAgent().review(deck)
        ok1 = reviewed.slides[0].points == ["Phần 1: Giới thiệu", "Phần 2: Phân tích"]
        ok2 = reviewed.slides[1].points == ["**Ý thứ nhất** cần tối ưu"]
    finally:
        llm.generate = orig_gen
        
    ok = ok1 and ok2
    print(f"[9] Agenda Review Bypass: {'PASS' if ok else 'FAIL'} (slide1={reviewed.slides[0].points}, slide2={reviewed.slides[1].points})")
    return ok


if __name__ == "__main__":
    print("=== SELF-TEST (không cần Ollama) ===")
    a = test_parsers()
    b = test_pipeline_fallback()
    c = test_markdown_stripper()
    d = test_instruction_leaks()
    e = test_output_delimiter()
    f = test_title_formatting()
    g = test_deep_leakage_filter()
    h = test_deduplication()
    i = test_agenda_review_bypass()
    print("\nKẾT QUẢ:", "TẤT CẢ PASS ✅" if (a and b and c and d and e and f and g and h and i) else "CÓ LỖI ❌")
